# pdf2pptx.py
import os
import sys
import argparse
from pptx import Presentation
from pdf2image import convert_from_path
from tqdm import tqdm

def pdf_to_pptx(pdf_path, output_path, dpi=600):
    images = convert_from_path(pdf_path, dpi=dpi)

    prs = Presentation()

    for img in tqdm(images, desc="Converting PDF to PPTX", unit="page"):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        temp_image_path = "temp_page_image.png"
        img.save(temp_image_path, format="PNG")

        img_ratio = img.width / img.height
        slide_ratio = prs.slide_width / prs.slide_height

        if img_ratio > slide_ratio:
            width = prs.slide_width
            height = int(prs.slide_width / img_ratio)
        else:
            height = prs.slide_height
            width = int(prs.slide_height * img_ratio)

        slide.shapes.add_picture(temp_image_path, 0, 0, width=width, height=height)

        # Remove the temporary image file after adding it to the slide
        os.remove(temp_image_path)

    prs.save(output_path)
    # Print the output file path after saving
    print(f"Saved PPTX file to {output_path}")

def main():
    parser = argparse.ArgumentParser(description='Convert PDF to PPTX with specified DPI.')
    parser.add_argument('pdf_path', help='Path to the input PDF file.')
    parser.add_argument('--output_path', help='Path to the output PPTX file. If not provided, replaces .pdf with .pptx in the input file name.')
    parser.add_argument('--dpi', type=int, default=600, help='DPI for conversion, default is 600.')
    args = parser.parse_args()

    # If output_path is not provided, replace .pdf with .pptx in the input file name
    if args.output_path is None:
        base_name = os.path.splitext(args.pdf_path)[0]
        args.output_path = base_name + '.pptx'

    pdf_to_pptx(args.pdf_path, args.output_path, args.dpi)

if __name__ == "__main__":
    main()
